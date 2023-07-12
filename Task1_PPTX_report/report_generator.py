import collections.abc  # needed import for python 3.10
import argparse
import json
import pandas as pd

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData


# generates report based on a config file
class ReportGenerator:
    def __init__(self, config_file):
        self.config_file = config_file
        self.prs = Presentation()

    def generate_report(self):
        with open(self.config_file) as f:
            config = json.load(f)
        # iterate through the slides
        for slide in config["presentation"]:
            slide_type = slide["type"]
            slide_title = slide["title"]
            slide_content = slide["content"]

            if slide_type == "title":
                self.create_title_slide(slide_title, slide_content)
            elif slide_type == "text":
                self.create_text_slide(slide_title, slide_content)
            elif slide_type == "list":
                self.create_list_slide(slide_title, slide_content)
            elif slide_type == "picture":
                self.create_picture_slide(slide_title, slide_content)
            elif slide_type == "plot":
                configuration = slide["configuration"]
                self.create_plot_slide(slide_title, slide_content, configuration)

        self.prs.save("report.pptx")

    # Creates a title type slide
    def create_title_slide(self, title, content):
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
        slide.placeholders[1].text = content

    # Creates a text type slide
    def create_text_slide(self, title, content):
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
        left, top, width, height = Cm(3.5), Cm(3), Cm(14.5), Cm(30)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.text = content

    # Creates a list type slide
    def create_list_slide(self, title, content):
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
        content_placeholder = slide.shapes.placeholders[1].text_frame
        for item in content:
            p = content_placeholder.add_paragraph()
            p.level = item["level"]
            p.text = item["text"]

    # Creates a picture type slide
    def create_picture_slide(self, title, content):
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
        left, top = Cm(3.5), Cm(3)
        slide.shapes.add_picture(content, left, top)

    # Creates a plot type slide
    def create_plot_slide(self, title, content, configuration):
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
        data = pd.read_csv(
            content, sep=" ", header=0
        )  # change 'sep' when .dat delimiter is known
        chart_data = ChartData()
        chart_data.categories = data["x"]
        chart_data.add_series(None, data["y"])

        left, top, width, height = (
            Cm(3.5),
            Cm(3),
            Cm(16),
            Cm(12),
        )  # unknown info about width, height
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE, left, top, width, height, chart_data
        ).chart

        # find better ways for fonts
        xaxis = chart.category_axis
        yaxis = chart.value_axis
        yaxis.has_major_gridlines = yaxis.has_major_gridlines = False
        x_label, y_label = xaxis.axis_title.text_frame, yaxis.axis_title.text_frame
        p_x, p_y = x_label.paragraphs[0], y_label.paragraphs[0]
        run_x, run_y = p_x.add_run(), p_y.add_run()
        font_x, font_y = run_x.font, run_y.font
        run_x.text = configuration["x-label"]
        run_y.text = configuration["y-label"]
        font_x.size = font_y.size = Pt(12)
        font_x.bold = font_y.bold = False

        chart.has_legend = False


def main():
    # argument needed to pass the config file
    parser = argparse.ArgumentParser(
        description="Generate a report in pptx format based on a configuration file."
    )
    parser.add_argument("config_file", help="Path to the configuration file")
    args = parser.parse_args()

    report_generator = ReportGenerator(args.config_file)
    report_generator.generate_report()


if __name__ == "__main__":
    main()
